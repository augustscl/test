
from pptx import Presentation
from pptx.util import Inches
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)  # 16:9 宽屏
prs.slide_height = Inches(5.625)

# 图片目录
image_dir = "/Users/suchuanlei/.openclaw/workspace/slide-deck/从会聊到会干活/jpg/"

# 图片列表（按顺序）
image_files = [
    "01-slide-cover.jpg",
    "02-slide-yesterday.jpg",
    "03-slide-today.jpg",
    "04-slide-five-closures.jpg",
    "05-slide-requirement-closed-loop.jpg",
    "06-slide-skill-closed-loop.jpg",
    "07-slide-memory-closed-loop.jpg",
    "08-slide-execution-closed-loop.jpg",
    "09-slide-error-correction-closed-loop.jpg",
    "10-slide-case1-getnote.jpg",
    "11-slide-case2-voice.jpg",
    "12-slide-case3-wechat.jpg",
    "13-slide-summary.jpg",
    "14-slide-methodology.jpg",
    "15-slide-ending.jpg",
    "16-slide-thanks.jpg"
]

# 每张图片创建一个幻灯片
for img_file in image_files:
    img_path = os.path.join(image_dir, img_file)
    if os.path.exists(img_path):
        # 创建空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加图片，占满整个幻灯片
        left = top = Inches(0)
        height = prs.slide_height
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 保存PowerPoint文件
output_path = "/Users/suchuanlei/.openclaw/workspace/slide-deck/从会聊到会干活/从会聊到会干活-compressed.pptx"
prs.save(output_path)
print(f"PowerPoint文件已生成: {output_path}")

