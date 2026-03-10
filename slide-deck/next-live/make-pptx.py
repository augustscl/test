
from pptx import Presentation
from pptx.util import Inches
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)  # 16:9 宽屏
prs.slide_height = Inches(5.625)

# 图片目录
image_dir = "/Users/suchuanlei/.openclaw/workspace/slide-deck/next-live/images-google/"

# 图片列表（按顺序）
image_files = [
    "01-slide-cover.png",
    "02-slide-opening.png",
    "03-slide-five-closures.png",
    "04-slide-getnote-1.png",
    "05-slide-getnote-2.png",
    "06-slide-getnote-3.png",
    "07-slide-voice-1.png",
    "08-slide-voice-2.png",
    "09-slide-voice-3.png",
    "10-slide-wechat-1.png",
    "11-slide-wechat-2.png",
    "12-slide-wechat-3.png",
    "13-slide-summary.png",
    "14-slide-methodology.png",
    "15-slide-ending.png",
    "16-slide-thanks.png"
]

# 每张图片创建一个幻灯片
for img_file in image_files:
    img_path = os.path.join(image_dir, img_file)
    if os.path.exists(img_path):
        # 创建空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加图片，占满整个幻灯片
        left = top = Inches(0)
        height = prs.slide_height
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 保存PowerPoint文件
output_path = "/Users/suchuanlei/.openclaw/workspace/slide-deck/next-live/从会聊到会干活-Google版.pptx"
prs.save(output_path)
print(f"PowerPoint文件已生成: {output_path}")
